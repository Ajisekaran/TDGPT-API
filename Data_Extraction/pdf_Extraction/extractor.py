import os
import json
import io
import re
import time
import fitz  
import pandas as pd
import pytesseract
import aiofiles
from PIL import Image
from utils import summarize_text
from openpyxl import load_workbook
from markdown_it import MarkdownIt
from pptx import Presentation
from image_caption import describe_image


async def extract_pdf_content(file_path, groq_client, model):
    start = time.time()
    pdf = fitz.open(file_path)
    filename = os.path.splitext(os.path.basename(file_path))[0]

    page_outputs = []
    IMG_ROOT = os.path.join("output", "images")
    IMG_SUMMARY_DIR = os.path.join(IMG_ROOT, "img_summary")
    IMG_VISION_DIR = os.path.join(IMG_ROOT, "img_vision")
    os.makedirs(IMG_SUMMARY_DIR, exist_ok=True)
    os.makedirs(IMG_VISION_DIR, exist_ok=True)

    for i, page in enumerate(pdf, start=1):
        text = page.get_text("text").strip()
        images = []
        img_summary_files = []
        img_vision_files = []

        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list, start=1):
            xref = img[0]
            base_image = pdf.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]

            image_filename = f"{filename}_page{i}_img{img_index}.{image_ext}"
            image_path = os.path.join(IMG_ROOT, image_filename)

            async with aiofiles.open(image_path, "wb") as img_file:
                await img_file.write(image_bytes)
            images.append(image_path)

            image_obj = Image.open(io.BytesIO(image_bytes)).convert("RGB")
            ocr_text = pytesseract.image_to_string(image_obj).strip()

            if ocr_text:
                summary_path = os.path.join(IMG_SUMMARY_DIR, f"{filename}_page{i}_img{img_index}_summary.json")
                async with aiofiles.open(summary_path, "w", encoding="utf-8") as f:
                    await f.write(json.dumps({"ocr_text": ocr_text}, indent=2, ensure_ascii=False))
                img_summary_files.append(summary_path)
            else:
                vision_path = os.path.join(IMG_VISION_DIR, f"{filename}_page{i}_img{img_index}_vision.json")
                description = await describe_image(image_path, groq_client, model)
                async with aiofiles.open(vision_path, "w", encoding="utf-8") as f:
                    await f.write(json.dumps({"description": description}, indent=2, ensure_ascii=False))
                img_vision_files.append(vision_path)

        summary = await summarize_text(text, groq_client, model)

        page_outputs.append({
            "page_number": i,
            "text": text or "No text found.",
            "tables": "No table support in fitz. Consider tabula or pdfplumber.",
            "images": images,
            "img_summary_files": img_summary_files,
            "img_vision_files": img_vision_files,
            "summary": summary,
            "time_taken": f"{time.time() - start:.2f} sec"
        })

    return {
        "metadata": {
            "file_name": os.path.basename(file_path),
            "file_type": "pdf",
            "file_size": f"{os.path.getsize(file_path) / (1024 * 1024):.2f} MB",
            "page_count": len(pdf),
        },
        "pages": page_outputs,
        "overall_summary": "PDF extraction complete.",
        "total_time_taken": f"{time.time() - start:.2f} seconds"
    }


async def extract_excel_content(file_path, groq_client, model):
    content = []
    try:
        sheets = pd.read_excel(file_path, sheet_name=None)
    except Exception as e:
        return {"error": f"Failed to read Excel file: {str(e)}"}

    for sheet_name, df in sheets.items():
        df = df.fillna("").astype(str)

        for i, row in df.iterrows():
            row_data = row.to_dict()
            if any(cell.strip() for cell in row_data.values()):
                content.append({
                    "sheet": sheet_name,
                    "row_number": i + 1,
                    "row_data": row_data
                })

    return {
        "metadata": {
            "file_name": os.path.basename(file_path),
            "file_type": "excel",
            "file_size": f"{os.path.getsize(file_path) / 1024:.2f} KB",
            "sheet_count": len(sheets)
        },
        "rows_extracted": len(content),
        "content": content,
        "summary": "Excel data extracted successfully."
    }


async def extract_ppt_content(file_path, groq_client, model):
    start = time.time()
    prs = Presentation(file_path)
    filename = os.path.splitext(os.path.basename(file_path))[0]

    IMG_ROOT = os.path.join("output", "images")
    IMG_SUMMARY_DIR = os.path.join(IMG_ROOT, "img_summary")
    IMG_VISION_DIR = os.path.join(IMG_ROOT, "img_vision")
    os.makedirs(IMG_SUMMARY_DIR, exist_ok=True)
    os.makedirs(IMG_VISION_DIR, exist_ok=True)

    slide_outputs = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []
        image_files = []
        img_summary_files = []
        img_vision_files = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                image_filename = f"{filename}_slide{slide_num}_img.{image_ext}"
                image_path = os.path.join(IMG_ROOT, image_filename)

                async with aiofiles.open(image_path, "wb") as img_file:
                    await img_file.write(image_bytes)
                image_files.append(image_path)

                try:
                    img_obj = Image.open(io.BytesIO(image_bytes)).convert("RGB")
                    ocr_text = pytesseract.image_to_string(img_obj).strip()

                    if ocr_text:
                        summary_path = os.path.join(IMG_SUMMARY_DIR, f"{filename}_slide{slide_num}_summary.json")
                        async with aiofiles.open(summary_path, "w", encoding="utf-8") as f:
                            await f.write(json.dumps({"ocr_text": ocr_text}, indent=2, ensure_ascii=False))
                        img_summary_files.append(summary_path)
                    else:
                        description = await describe_image(image_path, groq_client, model)
                        vision_path = os.path.join(IMG_VISION_DIR, f"{filename}_slide{slide_num}_vision.json")
                        async with aiofiles.open(vision_path, "w", encoding="utf-8") as f:
                            await f.write(json.dumps({"description": description}, indent=2, ensure_ascii=False))
                        img_vision_files.append(vision_path)

                except Exception as e:
                    print(f"Error processing image on slide {slide_num}: {e}")

        combined_text = "\n".join(slide_text)
        summary = await summarize_text(combined_text, groq_client, model) if combined_text else "No text to summarize."

        slide_outputs.append({
            "slide_number": slide_num,
            "text_blocks": slide_text,
            "images": image_files,
            "img_summary_files": img_summary_files,
            "img_vision_files": img_vision_files,
            "summary": summary
        })

    return {
        "metadata": {
            "file_name": os.path.basename(file_path),
            "file_type": "pptx",
            "file_size": f"{os.path.getsize(file_path) / (1024 * 1024):.2f} MB",
            "slide_count": len(prs.slides)
        },
        "slides": slide_outputs,
        "overall_summary": "PPTX extraction complete.",
        "total_time_taken": f"{time.time() - start:.2f} seconds"
    }


async def extract_markdown_content(file_path, groq_client=None, model=None):
    async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
        markdown_text = await f.read()

    file_size_kb = round(os.path.getsize(file_path) / 1024, 2)
    md = MarkdownIt()
    tokens = md.parse(markdown_text)

    structured_data = {
        "metadata": {
            "file_name": os.path.basename(file_path),
            "file_type": "markdown",
            "file_size": f"{file_size_kb} KB"
        },
        "content": []
    }

    current_block = {}
    current_type = None

    for token in tokens:
        if token.type == "heading_open":
            current_type = "heading"
        elif token.type == "paragraph_open":
            current_type = "paragraph"
        elif token.type == "bullet_list_open":
            current_type = "list"
            current_block = {"type": "list", "items": []}
        elif token.type == "list_item_open":
            current_type = "list_item"
        elif token.type == "inline":
            content = token.content.strip()
            if current_type == "heading":
                structured_data["content"].append({"type": "heading", "text": content})
            elif current_type == "paragraph":
                structured_data["content"].append({"type": "paragraph", "text": content})
            elif current_type == "list_item":
                current_block["items"].append(content)
        elif token.type == "bullet_list_close":
            structured_data["content"].append(current_block)
            current_block = {}
        elif token.type == "fence":
            structured_data["content"].append({
                "type": "code_block",
                "language": token.info.strip(),
                "code": token.content
            })

    return structured_data
